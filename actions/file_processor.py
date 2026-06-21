"""file_processor.py — Universal file processor for images, PDFs, docs, code, audio, video."""
from __future__ import annotations
import os, subprocess, sys, json
from pathlib import Path


def file_processor(parameters=None, player=None, speak=None, **kwargs):
    if parameters is None:
        parameters = {}
    file_path = parameters.get("file_path", "")
    action = parameters.get("action", "")
    instruction = parameters.get("instruction", "")
    fmt = parameters.get("format", "")
    dest = parameters.get("destination", "")

    if not file_path or not os.path.isfile(file_path):
        return "Archivo no encontrado. Proporcioná file_path válido."

    ext = Path(file_path).suffix.lower()
    fname = Path(file_path).name

    try:
        # ── Images ──────────────────────────────────────────────────────────
        if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"):
            from PIL import Image
            img = Image.open(file_path)
            if action == "describe":
                return f"Imagen {img.size}, modo {img.mode}, formato {img.format}. Usá screen_vision para analizar."
            elif action == "ocr":
                try:
                    import pytesseract
                    text = pytesseract.image_to_string(img)
                    return text[:1500] if text.strip() else "No se detectó texto."
                except ImportError:
                    return "pytesseract no instalado."
            elif action == "resize":
                w = parameters.get("width", 800)
                h = parameters.get("height", 600)
                scale = parameters.get("scale", 1)
                if scale != 1:
                    w = int(img.width * scale)
                    h = int(img.height * scale)
                img = img.resize((w, h), Image.LANCZOS)
                out = dest or file_path
                img.save(out)
                return f"Imagen redimensionada a {w}x{h}: {out}"
            elif action == "compress":
                quality = parameters.get("quality", 85)
                out = dest or file_path
                img.save(out, optimize=True, quality=quality)
                return f"Imagen comprimida (quality={quality}): {out}"
            elif action == "convert":
                out = dest or str(Path(file_path).with_suffix(f".{fmt or 'png'}"))
                img.save(out)
                return f"Imagen convertida: {out}"
            elif action == "info":
                return f"Info: {fname}\nTamaño: {img.size}\nModo: {img.mode}\nFormato: {img.format}\nPeso: {os.path.getsize(file_path)} bytes"

        # ── PDFs ────────────────────────────────────────────────────────────
        elif ext == ".pdf":
            try:
                import PyPDF2
            except ImportError:
                return "PyPDF2 no instalado. Ejecutá: pip install PyPDF2"
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                if action == "summarize":
                    text = " ".join(p.extract_text() for p in reader.pages[:5])
                    return f"Resumen ({len(reader.pages)} págs):\n{text[:1500]}"
                elif action == "extract_text":
                    text = " ".join(p.extract_text() for p in reader.pages)
                    return text[:3000] if text else "No se pudo extraer texto."
                elif action == "info":
                    return f"Info PDF: {fname}\nPáginas: {len(reader.pages)}"

        # ── DOCX / TXT ──────────────────────────────────────────────────────
        elif ext in (".txt", ".md", ".py", ".js", ".ts", ".html", ".css", ".json", ".xml", ".yaml", ".yml", ".csv"):
            content = Path(file_path).read_text(encoding="utf-8", errors="ignore")
            if action == "summarize":
                lines = content.splitlines()
                return f"Resumen de {fname} ({len(lines)} líneas, {len(content)} chars):\n{content[:1000]}"
            elif action == "word_count":
                words = len(content.split())
                return f"{fname}: {len(content.splitlines())} líneas, {words} palabras, {len(content)} caracteres."
            elif action in ("fix", "reformat"):
                if instruction:
                    return f"Instrucción '{instruction}' recibida para {fname}. Usá code_helper action=edit para modificarlo."
                return f"{fname} leído ({len(content)} chars)."
            return content[:2000]

        # ── Excel ───────────────────────────────────────────────────────────
        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
            except ImportError:
                return "openpyxl no instalado."
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            info = []
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                info.append(f"  {ws.title}: {len(rows)} filas x {len(rows[0]) if rows else 0} cols")
            return f"Archivo Excel: {fname}\n" + "\n".join(info)

        # ── CSV ─────────────────────────────────────────────────────────────
        elif ext == ".csv":
            import csv
            with open(file_path, newline="", encoding="utf-8") as f:
                reader = csv.reader(f)
                rows = list(reader)
                header = rows[0] if rows else []
                return f"CSV: {fname}\nColumnas: {', '.join(header)}\nFilas: {len(rows)-1}"

        # ── PPTX ────────────────────────────────────────────────────────────
        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return "python-pptx no instalado."
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides[:5]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return f"Presentación: {fname} ({len(prs.slides)} diapositivas)\n{chr(10).join(texts[:20])[:2000]}"

        # ── Audio ───────────────────────────────────────────────────────────
        elif ext in (".mp3", ".wav", ".ogg", ".flac", ".m4a"):
            try:
                import speech_recognition as sr
            except ImportError:
                try:
                    import whisper
                    model = whisper.load_model("base")
                    result = model.transcribe(file_path)
                    return result["text"][:1500]
                except ImportError:
                    return "Procesamiento de audio: instalá whisper (pip install openai-whisper) o speechrecognition."
            r = sr.Recognizer()
            with sr.AudioFile(file_path) as source:
                audio = r.record(source)
            try:
                text = r.recognize_google(audio, language="es-ES")
                return f"Transcripción:\n{text[:1500]}"
            except Exception:
                return "No se pudo transcribir el audio."

        return f"Tipo de archivo '{ext}' no soportado aún para action '{action}'."
    except Exception as e:
        return f"Error procesando {fname}: {e}"
